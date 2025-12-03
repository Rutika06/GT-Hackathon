import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt


def _add_title_slide(prs: Presentation, dataset_name: str):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Automated Data Insight Report"

    subtitle = slide.placeholders[1]
    subtitle.text = f"Dataset: {dataset_name}"


def _add_insights_slide(prs: Presentation, insights_text: str):
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Executive Summary & Insights"

    body = slide.placeholders[1].text_frame
    body.clear()

    # Split insights into lines / bullets
    for line in insights_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        p = body.add_paragraph()
        p.text = line
        p.level = 0


def _add_chart_slide(prs: Presentation, df: pd.DataFrame, output_dir: str = "report"):
    # Pick numeric column to visualize
    numeric_cols = df.select_dtypes(include="number").columns
    if len(numeric_cols) == 0:
        return  # Nothing to chart

    col = numeric_cols[0]  # first numeric column

    # Create simple bar chart for first 10 rows
    fig, ax = plt.subplots()
    df[col].head(10).plot(kind="bar", ax=ax)
    ax.set_title(f"Top 10 values of '{col}'")
    ax.set_ylabel(col)
    plt.tight_layout()

    os.makedirs(output_dir, exist_ok=True)
    chart_path = os.path.join(output_dir, "temp_chart.png")
    fig.savefig(chart_path)
    plt.close(fig)

    # Add slide with chart
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title if slide.shapes.title else None
    if title_shape:
        title_shape.text = f"Chart: {col}"

    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(chart_path, left, top, width=Inches(8))


def generate_ppt_report(
    dataset_name: str,
    df: pd.DataFrame,
    insights_text: str,
    output_path: str,
) -> str:
    """
    Create a PPTX report with:
    - Title slide
    - Executive summary slide (LLM insights)
    - Simple bar chart slide for first numeric column
    """
    prs = Presentation()

    _add_title_slide(prs, dataset_name)
    _add_insights_slide(prs, insights_text)
    _add_chart_slide(prs, df, output_dir=os.path.dirname(output_path) or "report")

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    return output_path
